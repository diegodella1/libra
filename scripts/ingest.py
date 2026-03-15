"""
Libra — Ingesta de documentos desde ZIP.

Procesa archivos directamente desde un ZIP sin descomprimir todo,
uno a la vez, para funcionar en la Pi 5 con espacio limitado.

Uso:
  source venv/bin/activate
  python scripts/ingest.py --zip /path/to/DATIP.zip --db-password <password>
  python scripts/ingest.py --zip /path/to/DATIP.zip --db-password <password> --dry-run
  python scripts/ingest.py --zip /path/to/DATIP.zip --db-password <password> --skip-pdf
"""

import os
import sys
import zipfile
import argparse
import tempfile
import shutil
from pathlib import Path
from datetime import datetime

from lib.common import (
    KNOWN_PERSONS, EXCLUDE_PATTERNS, ROMAN_MAP,
    should_exclude, parse_roman, infer_doc_type_from_path,
    extract_platform_from_path, extract_person_from_path,
    match_person_in_path, parse_metadata_from_path,
    get_db_connection, get_person_id_cache, link_persons,
    log_error, get_processed_files,
)

# Extensiones soportadas
SUPPORTED_EXTENSIONS = {'.txt', '.docx', '.xlsx', '.pptx', '.pdf'}

# Prioridad de procesamiento (mas rapidos primero)
EXTENSION_PRIORITY = {'.txt': 0, '.docx': 1, '.xlsx': 2, '.pptx': 3, '.pdf': 4}

# Tamano maximo de PDF para procesar completo (50MB)
PDF_MAX_SIZE = 50 * 1024 * 1024
PDF_MAX_PAGES = 50

# Minimo de texto extraido para considerar un PDF como "con texto"
MIN_TEXT_LENGTH = 50

# Mapeo extension -> doc_type (fallback, se sobreescribe con inferencia de ruta)
DOC_TYPE_MAP = {
    '.txt': 'texto',
    '.docx': 'documento',
    '.xlsx': 'planilla',
    '.pptx': 'presentacion',
    '.pdf': 'pdf',
}

# Directorio destino para PDFs servidos por nginx
PDF_OUTPUT_DIR = '/home/diego/Documents/libra/documents'


def extract_text_txt(file_path: str) -> str:
    """Extrae texto de un archivo TXT."""
    encodings = ['utf-8', 'latin-1', 'cp1252']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    return ''


def extract_text_docx(file_path: str) -> str:
    """Extrae texto de un archivo DOCX."""
    from docx import Document
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return '\n\n'.join(paragraphs)


def extract_text_xlsx(file_path: str) -> str:
    """Extrae texto de un archivo XLSX."""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"=== Hoja: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            line = '\t'.join(cells).strip()
            if line:
                lines.append(line)
    wb.close()
    return '\n'.join(lines)


def extract_text_pptx(file_path: str) -> str:
    """Extrae texto de un archivo PPTX."""
    from pptx import Presentation
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)
        if len(slide_lines) > 1:  # Mas que solo el header
            slides_text.append('\n'.join(slide_lines))
    return '\n\n'.join(slides_text)


def extract_text_pdf(file_path: str, file_size: int) -> tuple[str, int, str]:
    """
    Extrae texto de un PDF.
    Returns: (text, page_count, ocr_status)
    """
    from PyPDF2 import PdfReader

    reader = PdfReader(file_path)
    total_pages = len(reader.pages)

    # Limitar paginas para PDFs grandes
    max_pages = PDF_MAX_PAGES if file_size > PDF_MAX_SIZE else total_pages
    pages_to_read = min(total_pages, max_pages)

    pages_text = []
    for i in range(pages_to_read):
        try:
            text = reader.pages[i].extract_text() or ''
            pages_text.append(text)
        except Exception:
            continue

    full_text = '\n\n'.join(pages_text)

    # Determinar si necesita OCR
    if len(full_text.strip()) < MIN_TEXT_LENGTH:
        ocr_status = 'pending'
    else:
        ocr_status = 'skipped'  # No necesita OCR, tiene texto nativo

    return full_text, total_pages, ocr_status


EXTRACTORS = {
    '.txt': lambda path, size: (extract_text_txt(path), None, 'skipped'),
    '.docx': lambda path, size: (extract_text_docx(path), None, 'skipped'),
    '.xlsx': lambda path, size: (extract_text_xlsx(path), None, 'skipped'),
    '.pptx': lambda path, size: (extract_text_pptx(path), None, 'skipped'),
    '.pdf': lambda path, size: extract_text_pdf(path, size),
}


def upsert_document(conn, file_path: str, metadata: dict, content: str,
                    file_size: int, page_count: int | None, ocr_status: str,
                    doc_type: str) -> str | None:
    """Inserta o actualiza un documento y su registro en ingestion_log. Devuelve document_id."""
    doc_id = None
    with conn.cursor() as cur:
        # Construir tags desde metadata
        tags = []
        if metadata.get('efecto'):
            tags.append(metadata['efecto'])
        if metadata.get('punto'):
            tags.append(metadata['punto'])
        if metadata.get('platform'):
            tags.append(metadata['platform'])

        # Upsert en documents
        cur.execute("""
            INSERT INTO documents (title, doc_type, date, participants, tags, content,
                                   file_path, file_size, page_count, ocr_status)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
            ON CONFLICT (file_path) DO UPDATE SET
                title = EXCLUDED.title,
                content = EXCLUDED.content,
                file_size = EXCLUDED.file_size,
                page_count = EXCLUDED.page_count,
                ocr_status = EXCLUDED.ocr_status,
                updated_at = now()
            RETURNING id
        """, (
            metadata['title'],
            doc_type,
            metadata.get('date'),
            metadata.get('participants', []),
            tags,
            content if content else None,
            file_path,
            file_size,
            page_count,
            ocr_status,
        ))
        row = cur.fetchone()
        if row:
            doc_id = row[0]

        # Upsert en ingestion_log
        status = 'ocr' if ocr_status == 'pending' else 'indexed'
        cur.execute("""
            INSERT INTO ingestion_log (file_path, status, completed_at)
            VALUES (%s, %s, now())
            ON CONFLICT (file_path) DO UPDATE SET
                status = EXCLUDED.status,
                completed_at = now()
        """, (file_path, status))

    conn.commit()
    return doc_id


def copy_pdf_to_output(zip_file, zip_info, output_dir: str):
    """Copia un PDF del ZIP al directorio de salida, manteniendo estructura simplificada."""
    # Sanitizar path: quitar prefijo ANEXOS/ y caracteres problematicos
    rel_path = zip_info.filename
    # Quitar prefijo comun si existe
    for prefix in ['ANEXOS/', 'anexos/']:
        if rel_path.startswith(prefix):
            rel_path = rel_path[len(prefix):]
            break

    dest_path = os.path.join(output_dir, rel_path)
    dest_dir = os.path.dirname(dest_path)
    os.makedirs(dest_dir, exist_ok=True)

    with zip_file.open(zip_info.filename) as src, open(dest_path, 'wb') as dst:
        shutil.copyfileobj(src, dst)

    return rel_path


def process_zip(zip_path: str, conn, dry_run: bool = False, skip_pdf: bool = False,
                copy_pdfs: bool = True):
    """Procesa todos los archivos soportados de un ZIP."""
    tmp_dir = '/tmp/libra_ingest'
    os.makedirs(tmp_dir, exist_ok=True)

    with zipfile.ZipFile(zip_path, 'r') as zf:
        # Filtrar y ordenar archivos
        entries = []
        for info in zf.infolist():
            if info.is_dir():
                continue
            if info.file_size == 0:
                continue
            if should_exclude(info.filename):
                continue

            ext = Path(info.filename).suffix.lower()
            if ext not in SUPPORTED_EXTENSIONS:
                continue
            if skip_pdf and ext == '.pdf':
                continue

            entries.append(info)

        # Ordenar: por prioridad de extension, luego por tamano
        entries.sort(key=lambda e: (
            EXTENSION_PRIORITY.get(Path(e.filename).suffix.lower(), 99),
            e.file_size
        ))

        total = len(entries)
        print(f"\nArchivos a procesar: {total}")

        # Contar por tipo
        by_type = {}
        for e in entries:
            ext = Path(e.filename).suffix.lower()
            by_type[ext] = by_type.get(ext, 0) + 1
        for ext, count in sorted(by_type.items()):
            print(f"  {ext}: {count}")

        if dry_run:
            print("\n--- DRY RUN: listando archivos ---")
            for e in entries:
                size_mb = e.file_size / (1024 * 1024)
                doc_type = infer_doc_type_from_path(e.filename) or DOC_TYPE_MAP.get(Path(e.filename).suffix.lower(), '?')
                persons = match_person_in_path(e.filename)
                person_str = f" -> {', '.join(persons)}" if persons else ""
                print(f"  [{doc_type:13s}] {size_mb:7.2f} MB  {e.filename}{person_str}")
            return

        # Cargar cache de personas
        person_cache = get_person_id_cache(conn)

        # Obtener archivos ya procesados
        processed = get_processed_files(conn)
        pending = [e for e in entries if e.filename not in processed]
        print(f"Ya procesados: {total - len(pending)}, pendientes: {len(pending)}")

        if not pending:
            print("Nada que procesar.")
            return

        # Crear directorio de PDFs si corresponde
        if copy_pdfs:
            os.makedirs(PDF_OUTPUT_DIR, exist_ok=True)

        # Procesar
        success = 0
        errors = 0
        skipped_ocr = 0
        pdfs_copied = 0

        for i, info in enumerate(pending, 1):
            ext = Path(info.filename).suffix.lower()

            try:
                # Extraer a temporal
                tmp_path = os.path.join(tmp_dir, os.path.basename(info.filename))
                with zf.open(info.filename) as src, open(tmp_path, 'wb') as dst:
                    shutil.copyfileobj(src, dst)

                # Parsear metadata
                metadata = parse_metadata_from_path(info.filename)

                # Extraer texto
                extractor = EXTRACTORS[ext]
                content, page_count, ocr_status = extractor(tmp_path, info.file_size)

                # Doc type: inferir de ruta, fallback a extension
                doc_type = infer_doc_type_from_path(info.filename) or DOC_TYPE_MAP.get(ext, 'otro')

                # Upsert
                doc_id = upsert_document(
                    conn=conn,
                    file_path=info.filename,
                    metadata=metadata,
                    content=content,
                    file_size=info.file_size,
                    page_count=page_count,
                    ocr_status=ocr_status,
                    doc_type=doc_type,
                )

                # Vincular personas
                matched_persons = match_person_in_path(info.filename)
                if matched_persons and doc_id:
                    link_persons(conn, doc_id, matched_persons, person_cache)

                # Copiar PDF al directorio de nginx
                if copy_pdfs and ext == '.pdf':
                    try:
                        copy_pdf_to_output(zf, info, PDF_OUTPUT_DIR)
                        pdfs_copied += 1
                    except Exception as e_copy:
                        print(f"  WARN: no se pudo copiar PDF {info.filename}: {e_copy}", file=sys.stderr)

                success += 1
                if ocr_status == 'pending':
                    skipped_ocr += 1

                # Progress cada 50 archivos o al final
                if i % 50 == 0 or i == len(pending):
                    print(f"  [{i}/{len(pending)}] OK: {success} | Err: {errors} | OCR pending: {skipped_ocr} | PDFs: {pdfs_copied} | Ultimo: {info.filename}")

            except Exception as e:
                errors += 1
                error_msg = f"{type(e).__name__}: {str(e)}"
                try:
                    log_error(conn, info.filename, error_msg)
                except Exception:
                    pass
                if errors <= 10:
                    print(f"  ERROR [{i}/{len(pending)}] {info.filename}: {error_msg}", file=sys.stderr)
                elif errors == 11:
                    print("  ... suprimiendo errores siguientes", file=sys.stderr)

            finally:
                # Limpiar archivo temporal
                try:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)
                except Exception:
                    pass

        print(f"\n=== RESUMEN ===")
        print(f"Procesados: {success}")
        print(f"Errores: {errors}")
        print(f"OCR pendiente: {skipped_ocr}")
        if copy_pdfs:
            print(f"PDFs copiados: {pdfs_copied}")

    # Limpiar directorio temporal
    try:
        shutil.rmtree(tmp_dir, ignore_errors=True)
    except Exception:
        pass


def main():
    parser = argparse.ArgumentParser(description='Libra — Ingesta de documentos desde ZIP')
    parser.add_argument('--zip', required=True, help='Path al archivo ZIP')
    parser.add_argument('--db-password', required=True, help='Password de supabase_admin')
    parser.add_argument('--db-host', default='192.168.1.14', help='Host de PostgreSQL')
    parser.add_argument('--db-port', type=int, default=54329, help='Puerto de PostgreSQL (directo, sin pooler)')
    parser.add_argument('--dry-run', action='store_true', help='Solo listar archivos, no procesar')
    parser.add_argument('--skip-pdf', action='store_true', help='Saltar archivos PDF')
    parser.add_argument('--no-copy-pdfs', action='store_true', help='No copiar PDFs al directorio de nginx')
    args = parser.parse_args()

    # Validar zip
    if not os.path.exists(args.zip):
        print(f"Error: archivo ZIP no encontrado: {args.zip}")
        sys.exit(1)

    if not zipfile.is_zipfile(args.zip):
        print(f"Error: no es un archivo ZIP valido: {args.zip}")
        sys.exit(1)

    conn = None
    if not args.dry_run:
        try:
            conn = get_db_connection(args.db_password, args.db_host, args.db_port)
            print(f"Conectado a PostgreSQL en {args.db_host}:{args.db_port}")
        except Exception as e:
            print(f"Error conectando a DB: {e}")
            sys.exit(1)

    try:
        process_zip(args.zip, conn, args.dry_run, args.skip_pdf, not args.no_copy_pdfs)
    finally:
        if conn:
            conn.close()


if __name__ == '__main__':
    main()
